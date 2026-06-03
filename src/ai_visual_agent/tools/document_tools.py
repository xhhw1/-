from pathlib import Path
from typing import Any
import mimetypes

from langchain_core.tools import tool

from ai_visual_agent.config import get_settings


def _resolve_file_uri(file_uri: str) -> Path:
    if file_uri.startswith("file://"):
        file_uri = file_uri.removeprefix("file://")
    return Path(file_uri).expanduser().resolve()


def _summarize_text(text: str, max_chars: int = 220) -> str:
    compact = " ".join(text.split())
    if len(compact) <= max_chars:
        return compact
    return f"{compact[:max_chars].rstrip()}..."


def _parse_pptx(path: Path) -> list[dict[str, Any]]:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover - optional dependency guard
        raise RuntimeError("python-pptx is required to parse PPTX files.") from exc

    presentation = Presentation(str(path))
    pages: list[dict[str, Any]] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        text_parts: list[str] = []
        image_ids: list[str] = []

        title = ""
        if slide.shapes.title and getattr(slide.shapes.title, "text", None):
            title = slide.shapes.title.text.strip()

        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text:
                text_parts.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        text_parts.append(" | ".join(cells))
            if getattr(shape, "shape_type", None) and "PICTURE" in str(shape.shape_type):
                image_ids.append(f"{path.stem}-slide-{idx}-image-{len(image_ids) + 1}")

        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    text_parts.append(f"Notes: {notes}")
        except Exception:
            pass

        text = "\n".join(part for part in text_parts if part)
        pages.append(
            {
                "page_index": idx,
                "title": title or None,
                "text": text,
                "ocr_text": "",
                "image_asset_ids": image_ids,
                "semantic_summary": _summarize_text(text),
            }
        )
    return pages


def _parse_pdf(path: Path) -> list[dict[str, Any]]:
    try:
        import fitz
    except ImportError as exc:  # pragma: no cover - optional dependency guard
        raise RuntimeError("PyMuPDF is required to parse PDF files.") from exc

    pages: list[dict[str, Any]] = []
    with fitz.open(path) as document:
        for idx, page in enumerate(document, start=1):
            text = page.get_text("text").strip()
            first_line = next((line.strip() for line in text.splitlines() if line.strip()), "")
            image_ids = [
                f"{path.stem}-page-{idx}-image-{image_index}"
                for image_index, _image in enumerate(page.get_images(full=True), start=1)
            ]
            pages.append(
                {
                    "page_index": idx,
                    "title": first_line[:80] or None,
                    "text": text,
                    "ocr_text": "",
                    "image_asset_ids": image_ids,
                    "semantic_summary": _summarize_text(text),
                }
            )
    return pages


def _parse_llamaparse(path: Path) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    settings = get_settings()
    if not settings.llama_cloud_api_key:
        raise RuntimeError("LLAMA_CLOUD_API_KEY is required for DOCUMENT_PARSER_BACKEND=llamaparse.")

    try:
        from llama_cloud import LlamaCloud
    except ImportError as exc:  # pragma: no cover - optional dependency guard
        raise RuntimeError("DOCUMENT_PARSER_BACKEND=llamaparse requires llama-cloud.") from exc

    mime_type = mimetypes.guess_type(path.name)[0] or "application/octet-stream"
    client = LlamaCloud(api_key=settings.llama_cloud_api_key)
    response = client.parsing.parse(
        tier=settings.llama_parse_tier,  # type: ignore[arg-type]
        version=settings.llama_parse_version,
        expand=_llamaparse_expand_fields(settings.llama_parse_tier),
        upload_file=(path.name, path.read_bytes(), mime_type),
        timeout=settings.llama_parse_timeout,
    )
    return _pages_from_llamaparse_response(response), _llamaparse_metadata(response)


def _llamaparse_expand_fields(tier: str) -> list[str]:
    if tier.lower() == "fast":
        return ["text", "items", "job_metadata"]
    return ["markdown", "text", "items", "job_metadata"]


def _pages_from_llamaparse_response(response: Any) -> list[dict[str, Any]]:
    items = _object_get(response, "items") or []
    pages: list[dict[str, Any]] = []
    if isinstance(items, list) and items:
        for idx, item in enumerate(items, start=1):
            text = _string_value(
                _object_get(item, "text")
                or _object_get(item, "markdown")
                or _object_get(item, "content")
            )
            title = next((line.strip("# ").strip() for line in text.splitlines() if line.strip()), "")
            pages.append(
                {
                    "page_index": int(_object_get(item, "page") or _object_get(item, "page_index") or idx),
                    "title": title[:80] or None,
                    "text": text,
                    "ocr_text": "",
                    "image_asset_ids": _string_list(_object_get(item, "image_asset_ids")),
                    "semantic_summary": _summarize_text(text),
                }
            )
        return pages

    text_full = _string_value(
        _object_get(response, "text_full")
        or _object_get(response, "markdown_full")
        or _object_get(response, "text")
        or _object_get(response, "markdown")
    )
    chunks = [chunk.strip() for chunk in text_full.split("\n---\n") if chunk.strip()] or [text_full]
    for idx, text in enumerate(chunks, start=1):
        first_line = next((line.strip("# ").strip() for line in text.splitlines() if line.strip()), "")
        pages.append(
            {
                "page_index": idx,
                "title": first_line[:80] or None,
                "text": text,
                "ocr_text": "",
                "image_asset_ids": [],
                "semantic_summary": _summarize_text(text),
            }
        )
    return pages


def _llamaparse_metadata(response: Any) -> dict[str, Any]:
    job = _object_get(response, "job") or {}
    metadata = _object_get(response, "metadata") or {}
    return {
        "job_id": _object_get(job, "id") or _object_get(job, "job_id"),
        "status": _object_get(job, "status"),
        "metadata": metadata if isinstance(metadata, dict) else {},
    }


def _object_get(value: Any, key: str) -> Any:
    if isinstance(value, dict):
        return value.get(key)
    return getattr(value, key, None)


def _string_value(value: Any) -> str:
    if value is None:
        return ""
    pages = _object_get(value, "pages")
    if isinstance(pages, list):
        return "\n".join(
            _string_value(_object_get(page, "text") or _object_get(page, "markdown") or page)
            for page in pages
            if _string_value(_object_get(page, "text") or _object_get(page, "markdown") or page).strip()
        )
    if isinstance(value, list):
        return "\n".join(str(item) for item in value if str(item).strip())
    text = _object_get(value, "text")
    if text is not None and text is not value:
        return _string_value(text)
    markdown = _object_get(value, "markdown")
    if markdown is not None and markdown is not value:
        return _string_value(markdown)
    return str(value)


def _string_list(value: Any) -> list[str]:
    if value is None:
        return []
    if isinstance(value, list):
        return [str(item) for item in value if str(item).strip()]
    if isinstance(value, str):
        return [value] if value.strip() else []
    return [str(value)]


def parse_document_file(file_id: str, file_uri: str, file_type: str | None = None) -> dict[str, Any]:
    """Parse PPTX/PDF source material into a normalized page list."""

    if get_settings().mock_external_tools and file_uri == "mock":
        return {
            "file_id": file_id,
            "file_uri": file_uri,
            "file_type": file_type,
            "parser": "mock",
            "pages": [
                {
                    "page_index": 1,
                    "title": "产品资料概览",
                    "text": "待接入 python-pptx / LlamaParse / PyMuPDF。",
                    "ocr_text": "",
                    "image_asset_ids": [],
                    "semantic_summary": "当前为 mock 输出，用于联通 LangGraph。",
                }
            ],
        }

    path = _resolve_file_uri(file_uri)
    if not path.exists():
        raise FileNotFoundError(f"Document not found: {path}")

    extension = path.suffix.lower()
    normalized_type = (file_type or extension).lower()
    parser_metadata: dict[str, Any] = {}
    parser_backend = get_settings().document_parser_backend.lower()
    if parser_backend in {"llamaparse", "llama_parse", "llama-cloud", "llama_cloud"}:
        try:
            pages, parser_metadata = _parse_llamaparse(path)
            parser = "llamaparse"
        except Exception as exc:
            if extension == ".pptx" or "presentation" in normalized_type:
                pages = _parse_pptx(path)
                parser = "python-pptx-fallback"
            elif extension == ".pdf" or "pdf" in normalized_type:
                pages = _parse_pdf(path)
                parser = "pymupdf-fallback"
            else:
                raise
            parser_metadata = {"primary_parser": "llamaparse", "primary_error": f"{type(exc).__name__}: {exc}"}
    elif extension == ".pptx" or "presentation" in normalized_type:
        pages = _parse_pptx(path)
        parser = "python-pptx"
    elif extension == ".pdf" or "pdf" in normalized_type:
        pages = _parse_pdf(path)
        parser = "pymupdf"
    else:
        raise ValueError(f"Unsupported document type: {file_type or extension}")

    return {
        "file_id": file_id,
        "file_uri": str(path),
        "file_type": file_type or extension,
        "parser": parser,
        "pages": pages,
        "parser_metadata": parser_metadata,
    }


@tool
def parse_document(file_id: str, file_uri: str, file_type: str | None = None) -> dict[str, Any]:
    """Parse PPT/PDF source material into page text, OCR text, images, and semantic summaries."""

    return parse_document_file(file_id=file_id, file_uri=file_uri, file_type=file_type)
