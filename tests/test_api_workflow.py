import importlib.util

import pytest
from fastapi.testclient import TestClient

from ai_visual_agent.main import app


pytestmark = pytest.mark.skipif(
    importlib.util.find_spec("pptx") is None,
    reason="python-pptx optional dependency is not installed",
)


def test_upload_pptx_and_start_workflow(tmp_path) -> None:
    from pptx import Presentation

    pptx_path = tmp_path / "product.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "产品资料"
    slide.shapes.add_textbox(0, 0, 5000000, 1000000).text = "尺寸：20cm\n配件：主体\n玩法：互动"
    presentation.save(pptx_path)

    client = TestClient(app)
    project_response = client.post(
        "/api/projects",
        json={
            "workflow_type": "packaging",
            "brief": {
                "category": "玩具",
                "target_user": "亲子家庭",
                "user_expectations": ["安全", "好玩"],
                "value_proposition": "更强互动体验",
                "core_product_definition": "互动玩具套装",
            },
            "assets": [],
        },
    )
    assert project_response.status_code == 200
    project_id = project_response.json()["id"]

    with pptx_path.open("rb") as handle:
        upload_response = client.post(
            f"/api/projects/{project_id}/assets",
            data={"kind": "product_ppt"},
            files={
                "file": (
                    "product.pptx",
                    handle,
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            },
        )

    assert upload_response.status_code == 200
    assert upload_response.json()["kind"] == "product_ppt"

    start_response = client.post(f"/api/workflows/{project_id}/start")
    assert start_response.status_code == 200
    body = start_response.json()
    assert body["status"] == "waiting_review"
    assert body["interrupts"][0]["value"]["type"] == "usp_review"
    parsed_product = body["state"]["parsed_product"]
    assert parsed_product["parsed_pages"][0]["title"] == "产品资料"
    assert parsed_product["dimensions"]
    assert body["state"]["memory_context"][0]["memory_ids"]

    audit_response = client.get(f"/api/projects/{project_id}/audit?record_type=agent_output")
    assert audit_response.status_code == 200
    assert {record["stage"] for record in audit_response.json()} >= {
        "parse_inputs",
        "analyze_competitors",
        "generate_usps",
    }

    first_resume = client.post(
        f"/api/workflows/{project_id}/resume",
        json={"action": "approve", "reviewer": "tester", "comment": "ok"},
    )
    assert first_resume.status_code == 200
    assert first_resume.json()["interrupts"][0]["value"]["type"] == "strategy_review"

    second_resume = client.post(
        f"/api/workflows/{project_id}/resume",
        json={"action": "approve", "reviewer": "tester", "comment": "ok"},
    )
    assert second_resume.status_code == 200
    assert second_resume.json()["interrupts"][0]["value"]["type"] == "final_design_review"

    final_resume = client.post(
        f"/api/workflows/{project_id}/resume",
        json={"action": "approve", "reviewer": "tester", "comment": "ok"},
    )
    assert final_resume.status_code == 200
    assert final_resume.json()["status"] == "completed"

    audit_records = client.get(f"/api/projects/{project_id}/audit").json()
    record_types = {record["record_type"] for record in audit_records}
    assert {"human_review", "agent_output", "qc_report", "archive_record"}.issubset(record_types)
