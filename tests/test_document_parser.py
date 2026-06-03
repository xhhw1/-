import importlib.util
from types import SimpleNamespace

import pytest

from ai_visual_agent.config import get_settings
from ai_visual_agent.tools import document_tools
from ai_visual_agent.tools.document_tools import parse_document_file


pytestmark = pytest.mark.skipif(
    importlib.util.find_spec("pptx") is None,
    reason="python-pptx optional dependency is not installed",
)


def test_parse_pptx_extracts_slide_text(monkeypatch, tmp_path) -> None:
    monkeypatch.setenv("DOCUMENT_PARSER_BACKEND", "local")
    get_settings.cache_clear()
    from pptx import Presentation

    pptx_path = tmp_path / "product.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "产品核心玩法"
    box = slide.shapes.add_textbox(0, 0, 5000000, 1000000)
    box.text = "尺寸：20cm x 10cm\n配件：主体、遥控器\n玩法：亲子互动"
    presentation.save(pptx_path)

    result = parse_document_file("asset-1", str(pptx_path), "product.pptx")

    assert result["parser"] == "python-pptx"
    assert result["pages"][0]["title"] == "产品核心玩法"
    assert "尺寸" in result["pages"][0]["text"]
    get_settings.cache_clear()


def test_llamaparse_response_is_normalized_to_pages() -> None:
    response = SimpleNamespace(
        items=[
            SimpleNamespace(page=1, markdown="# 首屏\n尺寸：20cm"),
            SimpleNamespace(page=2, text="配件：主体、遥控器"),
        ],
        job=SimpleNamespace(id="job-1", status="SUCCESS"),
        metadata={"source": "unit"},
    )

    pages = document_tools._pages_from_llamaparse_response(response)
    metadata = document_tools._llamaparse_metadata(response)

    assert pages[0]["title"] == "首屏"
    assert "配件" in pages[1]["text"]
    assert metadata["job_id"] == "job-1"
